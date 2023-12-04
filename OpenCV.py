import cv2
import mediapipe as mp
from pptx import Presentation
import pyautogui

# Initialize MediaPipe Hands
mp_hands = mp.solutions.hands
hands = mp_hands.Hands()

# Initialize PowerPoint
presentation = Presentation()

# Initialize OpenCV
cap = cv2.VideoCapture(0)
drawing = False
prev_point = (0, 0)

while cap.isOpened():
    # Read a frame from the webcam
    ret, frame = cap.read()
    if not ret:
        break

    # Convert the BGR image to RGB
    rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)

    # Process the frame with MediaPipe Hands
    results = hands.process(rgb_frame)

    # Draw hand landmarks on the frame and enable drawing
    if results.multi_hand_landmarks:
        for hand_landmarks in results.multi_hand_landmarks:
            for landmark in hand_landmarks.landmark:
                h, w, _ = frame.shape
                cx, cy = int(landmark.x * w), int(landmark.y * h)
                cv2.circle(frame, (cx, cy), 5, (0, 255, 0), -1)

            # Get the coordinates of the thumb tip
            thumb_tip = hand_landmarks.landmark[mp_hands.HandLandmark.THUMB_TIP]
            tx, ty = int(thumb_tip.x * w), int(thumb_tip.y * h)

            # Example: If thumb tip is on the right side, move to the next slide
            if tx > w / 2:
                presentation.slides.add_slide(presentation.slides[-1].slide_layout)

            # Example: If thumb tip is on the left side, move to the previous slide
            elif tx < w / 2:
                presentation.slides.add_slide(presentation.slides[0].slide_layout)

    # Display the frame
    cv2.imshow('Hand Gesture & PowerPoint Control', frame)

    # Check for key events
    key = cv2.waitKey(1)

    # Break the loop when 'q' is pressed
    if key & 0xFF == ord('q'):
        break

# Release the webcam and close all windows
cap.release()
cv2.destroyAllWindows()
